import os
import shutil
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def update_pptx(input_path, output_path):
    prs = Presentation(input_path)

    # Slide 1: Title Page
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Problem Statement ID" in text or "TITLE PAGE" in text or "Team ID" in text:
                shape.text_frame.clear()
                p1 = shape.text_frame.paragraphs[0]
                p1.text = "SMART INDIA HACKATHON 2025 / 2026"
                p1.font.bold = True
                p1.font.size = Pt(18)
                p1.font.color.rgb = RGBColor(0, 150, 255)

                p2 = shape.text_frame.add_paragraph()
                p2.text = "SIH25-806 • EV Charging Station Predictor"
                p2.font.bold = True
                p2.font.size = Pt(16)
                p2.font.color.rgb = RGBColor(15, 23, 42)

                details = [
                    ("Problem Statement Title", "EV Charging Station Predictor – Optimal Site Recommendation for Delhi NCT"),
                    ("Ministry / Agency", "Ministry of Housing & Urban Affairs (MoHUA)"),
                    ("Theme & Category", "Smart Vehicles / Transportation & Clean Tech (Software)"),
                    ("Team Name", "Hexagram"),
                    ("Team ID", "SIH2025-HEXAGRAM"),
                    ("Team Leader", "Sourav Kumar"),
                    ("Team Members", "Sumit Kumar, Rishabh Raj, Karishma Kumari Gupta, Abhas Prasad, Rishu Raj"),
                    ("Live Prototype URL", "https://souravkr114.github.io/modified-sih-project/")
                ]

                for label, val in details:
                    p = shape.text_frame.add_paragraph()
                    p.text = f"• {label}: {val}"
                    p.font.size = Pt(11)
                    p.font.color.rgb = RGBColor(30, 41, 59)

    # Slide 2: Proposed Solution
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_text_frame and "Proposed Solution" in shape.text_frame.text:
            shape.text_frame.clear()
            p_head = shape.text_frame.paragraphs[0]
            p_head.text = "PROPOSED SOLUTION: DELHI EV CHARGING INTELLIGENCE & INVESTMENT PLATFORM"
            p_head.font.bold = True
            p_head.font.size = Pt(13)
            p_head.font.color.rgb = RGBColor(2, 132, 199)

            pillars = [
                ("Pillar 1: Multi-Factor Spatial GIS Scoring (S_i ∈ [0, 100])",
                 "Fuses 6 weighted factors: EV fleet density (30%), supply gap (20%), grid substation capacity (15%), district adoption (15%), traffic corridors (10%), and POI dwell attraction (10%)."),
                ("Pillar 2: Explainable AI (SHAP Feature Rationale)",
                 "Replaces 'black-box ML' with game-theoretic additive feature attribution weights (φ_k) so city planners and CPOs see exactly WHY a location received Rank #1 (+24.5 traffic, +21.0 supply gap)."),
                ("Pillar 3: Investor Capex & Financial ROI Return Simulator",
                 "Calculates total Capex, DISCOM bulk power cost (₹6.8/kWh), annual ROI %, and payback period (months) across interactive tariff sliders (₹12–25/kWh)."),
                ("Dynamic Pincode & POI Engine",
                 "Enables real-time location scoring recalculation based on user Pincode, Search Radius (1–25km), and Points of Interest (Hotels, Malls, Metro Stations).")
            ]

            for title, desc in pillars:
                p_t = shape.text_frame.add_paragraph()
                p_t.text = f"• {title}:"
                p_t.font.bold = True
                p_t.font.size = Pt(10.5)
                p_t.font.color.rgb = RGBColor(15, 23, 42)

                p_d = shape.text_frame.add_paragraph()
                p_d.text = f"  {desc}"
                p_d.font.size = Pt(9.5)
                p_d.font.color.rgb = RGBColor(71, 85, 105)

    # Slide 3: Technical Approach
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.has_text_frame and "TECHNICAL APPROACH" in shape.text_frame.text:
            shape.text_frame.clear()
            p_head = shape.text_frame.paragraphs[0]
            p_head.text = "TECHNICAL APPROACH & PRODUCTION ARCHITECTURE"
            p_head.font.bold = True
            p_head.font.size = Pt(13)
            p_head.font.color.rgb = RGBColor(2, 132, 199)

            tech_items = [
                ("Frontend & GIS Mapping", "Next.js 16 (Turbopack), React 18, TypeScript, Tailwind CSS, Leaflet.js with CartoDB Voyager Light tiles & custom SVG rank markers."),
                ("Dynamic Predictor Engine", "Spatial Scoring Algorithm + Pincode, Radius (1–25km), and POI catchment recalculation module."),
                ("Explainable AI (XAI)", "SHAP (SHapley Additive exPlanations) game-theoretic feature decomposition engine."),
                ("Live & Offline Deployment Resilience", "Deployed live on GitHub Pages (souravkr114.github.io/modified-sih-project) + 100% offline single-file standalone bundle (Delhi_EV_Platform.html).")
            ]

            for title, desc in tech_items:
                p = shape.text_frame.add_paragraph()
                p.text = f"• {title}: {desc}"
                p.font.size = Pt(10.5)
                p.font.color.rgb = RGBColor(30, 41, 59)

    # Slide 4: Feasibility & Viability
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_text_frame and "FEASIBILITY AND VIABILITY" in shape.text_frame.text:
            shape.text_frame.clear()
            p_head = shape.text_frame.paragraphs[0]
            p_head.text = "FEASIBILITY, GRID CONSTRAINTS & DATA PROVENANCE"
            p_head.font.bold = True
            p_head.font.size = Pt(13)
            p_head.font.color.rgb = RGBColor(2, 132, 199)

            feasibility_items = [
                ("Data Provenance & Quality (Rule 2 Compliance)", "Ground truth Delhi Govt Vahan registrations & DISCOM substation capacities combined with OpenStreetMap GIS geometries and calibrated spatial proxies."),
                ("Grid Substation Feasibility Index", "Incorporates 11kV/33kV substation proximity, underground cable extension distance (meters), and transformer kVA readiness into location scoring."),
                ("Dynamic Supply Gap Adjustment", "Automatically recalculates surrounding deficit scores when new stations are commissioned to prevent market over-saturation."),
                ("Proven Cost-Effective Stack", "Zero server overhead using static Next.js production builds deployed to GitHub Pages.")
            ]

            for title, desc in feasibility_items:
                p = shape.text_frame.add_paragraph()
                p.text = f"• {title}: {desc}"
                p.font.size = Pt(10.5)
                p.font.color.rgb = RGBColor(30, 41, 59)

    # Slide 5: Impact & Benefits
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.has_text_frame and "IMPACT AND BENEFITS" in shape.text_frame.text:
            shape.text_frame.clear()
            p_head = shape.text_frame.paragraphs[0]
            p_head.text = "IMPACT, BENEFITS & MoHUA POLICY ALIGNMENT"
            p_head.font.bold = True
            p_head.font.size = Pt(13)
            p_head.font.color.rgb = RGBColor(2, 132, 199)

            impact_items = [
                ("MoHUA Urban Policy Alignment", "Direct tracking against official benchmark of 1 Public Charger per 25 EVs (highlighted in District Deficit Matrix)."),
                ("Economic Impact for Operators", "De-risks private CPO capital investments by predicting payback periods (14–24 months) before building."),
                ("Social & Environmental Impact", "Eliminates EV range anxiety, accelerates clean-air EV adoption across Delhi NCT, and prioritizes multi-modal transit hubs (ISBTs, Metro, Airport).")
            ]

            for title, desc in impact_items:
                p = shape.text_frame.add_paragraph()
                p.text = f"• {title}: {desc}"
                p.font.size = Pt(10.5)
                p.font.color.rgb = RGBColor(30, 41, 59)

    # Slide 6: Research & References
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_text_frame and "RESEARCH AND REFERENCES" in shape.text_frame.text:
            shape.text_frame.clear()
            p_head = shape.text_frame.paragraphs[0]
            p_head.text = "DEMO LINKS, CODE & REFERENCES"
            p_head.font.bold = True
            p_head.font.size = Pt(13)
            p_head.font.color.rgb = RGBColor(2, 132, 199)

            ref_items = [
                ("Live Interactive Prototype URL", "https://souravkr114.github.io/modified-sih-project/"),
                ("GitHub Code Repository", "https://github.com/souravkr114/modified-sih-project"),
                ("Policy References", "Delhi Electric Vehicle Policy 2026, MoHUA Urban EV Infrastructure Guidelines."),
                ("Spatial Data Sources", "OpenStreetMap (Roads, POIs), Delhi Metro Transit Network, DISCOM Substation Listings.")
            ]

            for title, desc in ref_items:
                p = shape.text_frame.add_paragraph()
                p.text = f"• {title}: {desc}"
                p.font.size = Pt(10.5)
                p.font.color.rgb = RGBColor(30, 41, 59)

    prs.save(output_path)
    print(f"Successfully updated presentation saved to: {output_path}")

if __name__ == "__main__":
    src = "unlocked_in.pptx"
    out = "SIH2026_Idea_Hexagram_DelhiEV.pptx"
    update_pptx(src, out)
