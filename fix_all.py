import os
import shutil
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def fix_powerpoint():
    src = 'SIH2025_Idea_DelhiEV.pptx'
    tmp = 'unlocked_temp.pptx'
    shutil.copy(src, tmp)

    prs = Presentation(tmp)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if 'SIH BV806' in r.text:
                            r.text = r.text.replace('SIH BV806', 'SIH25-806')
                        if 'BV806' in r.text:
                            r.text = r.text.replace('BV806', 'SIH25-806')
                        if '[Your Team ID]' in r.text:
                            r.text = r.text.replace('[Your Team ID]', 'SIH2025-HEXAGRAM')
                        if '[Your Team Name – as registered on portal]' in r.text or '[Your Team Name' in r.text:
                            r.text = r.text.replace('[Your Team Name – as registered on portal]', 'Hexagram').replace('[Your Team Name', 'Hexagram')

    # Specifically update Slide 1 details cleanly
    slide1 = prs.slides[0]
    # Check if we need to add member names shape cleanly
    found_members = False
    for shape in slide1.shapes:
        if shape.has_text_frame and 'Team Members' in shape.text_frame.text:
            found_members = True
            break
            
    if not found_members:
        # Add team details to Slide 1
        txBox = slide1.shapes.add_textbox(Pt(50), Pt(380), Pt(600), Pt(120))
        tf = txBox.text_frame
        p1 = tf.paragraphs[0]
        p1.text = "• Team Leader: Sourav Kumar"
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(15, 23, 42)

        p2 = tf.add_paragraph()
        p2.text = "• Members: Sumit Kumar, Rishabh Raj, Karishma Kumari Gupta, Abhas Prasad, Rishu Raj"
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(51, 65, 85)

        p3 = tf.add_paragraph()
        p3.text = "• Live Prototype: https://souravkr114.github.io/modified-sih-project/"
        p3.font.size = Pt(11)
        p3.font.bold = True
        p3.font.color.rgb = RGBColor(0, 150, 255)

    out_file = 'SIH2026_Idea_Hexagram_DelhiEV.pptx'
    prs.save(out_file)
    print(f"Fixed PowerPoint saved cleanly to: {out_file}")
    
    if os.path.exists(tmp):
        os.remove(tmp)

def fix_website_code():
    files_to_update = [
        'Delhi_EV_Platform.html',
        'index.html',
        'src/components/Navbar.tsx',
        'src/app/page.tsx',
        'src/components/DataLineageModal.tsx',
        'src/components/ReportExporter.tsx',
        'src/components/DistrictAnalytics.tsx'
    ]

    for fname in files_to_update:
        if os.path.exists(fname):
            with open(fname, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Replace all occurrences of BV806 with SIH25-806
            updated_content = content.replace('SIH BV806', 'SIH25-806').replace('BV806', 'SIH25-806')
            
            with open(fname, 'w', encoding='utf-8') as f:
                f.write(updated_content)
            print(f"Updated {fname} with SIH25-806")

if __name__ == "__main__":
    fix_powerpoint()
    fix_website_code()
